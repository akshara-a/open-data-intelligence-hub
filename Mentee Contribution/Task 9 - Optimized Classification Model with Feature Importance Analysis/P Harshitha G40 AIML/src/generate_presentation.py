import os
# pyrefly: ignore [missing-import]
from pptx import Presentation
# pyrefly: ignore [missing-import]
from pptx.util import Inches, Pt
# pyrefly: ignore [missing-import]
from pptx.enum.text import PP_ALIGN
# pyrefly: ignore [missing-import]
from pptx.dml.color import RGBColor

def create_presentation(output_pptx_path):
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6] # Blank slide

    # Theme colors
    PRIMARY_COLOR = RGBColor(26, 37, 48)     # Dark Slate
    ACCENT_COLOR = RGBColor(52, 152, 219)    # Bright Blue
    TEXT_DARK = RGBColor(44, 62, 80)         # Charcoal
    BG_LIGHT = RGBColor(245, 247, 250)       # Light Grey
    WHITE = RGBColor(255, 255, 255)

    slides_content = [
        {
            "title": "Predicting E-Commerce Purchase Likelihood",
            "subtitle": "Optimized Machine Learning Classification Model & Business Insights",
            "bullets": [
                "Objective: Predict whether a customer will complete a purchase based on browsing and demographic signals.",
                "Approach: End-to-end ML pipeline with Logistic Regression, Decision Trees, and Random Forest.",
                "Key Deliverables: Scikit-learn pipeline, high-precision classifier, threshold tuning, and actionable business strategies."
            ]
        },
        {
            "title": "1. Business Problem & Project Objective",
            "subtitle": "Maximizing Conversion Rates & Optimizing Marketing Budget",
            "bullets": [
                "E-Commerce Challenge: Low site-wide conversion rates lead to high customer acquisition costs and wasted promo spend.",
                "Core Question: 'Can we predict purchase intent from early session behavior and customer profiles?'",
                "Business Value:",
                "  • Target high-intent prospects with precision.",
                "  • Automate real-time cart recovery interventions.",
                "  • Segment low-intent traffic to reduce unnecessary ad spend."
            ]
        },
        {
            "title": "2. Dataset Overview & Data Quality",
            "subtitle": "1,500 Customer Session Records Across 18 Behavioral & Demographic Attributes",
            "bullets": [
                "Features Included: Age, Gender, Location, DeviceType, TrafficSource, PagesViewed, TimeOnSite, CartItems, PreviousPurchases, AverageOrderValue, DiscountUsed, EmailClicked, AdClicked, ReviewScoreViewed, DaysSinceLastVisit, SessionCount.",
                "Target Variable: Purchase (0 = Did Not Purchase: 70.8%, 1 = Purchased: 29.2%).",
                "Data Quality Checks: Handled missing values via median/mode imputation, removed duplicate records, and excluded CustomerID identifiers."
            ]
        },
        {
            "title": "3. Exploratory Data Analysis (EDA) Insights",
            "subtitle": "Key Behavioral & Channel Patterns Driving Conversions",
            "bullets": [
                "Cart Items Impact: Customers with >= 2 cart items achieve >85% purchase conversion rate.",
                "Engagement Depth: Time on site (>5 minutes) and pages viewed (>6 pages) double conversion odds.",
                "Repeat Buyers: Previous purchases strongly correlate with repeat buying behavior (3.2x higher conversion).",
                "Device & Channel: Desktop users exhibit higher conversion intent compared to mobile browsers."
            ]
        },
        {
            "title": "4. Preprocessing & Data Leakage Prevention",
            "subtitle": "Scikit-Learn ColumnTransformer Pipeline Architecture",
            "bullets": [
                "Preprocessing Steps:",
                "  • Numerical Features: SimpleImputer (median) + StandardScaler standardization.",
                "  • Categorical Features: SimpleImputer (most_frequent) + OneHotEncoder(handle_unknown='ignore').",
                "Data Leakage Prevention: Transformers fitted exclusively on training set within Scikit-Learn Pipelines.",
                "Stratified Data Split: 80% Training (1,200 rows) / 20% Testing (300 rows)."
            ]
        },
        {
            "title": "5. Baseline Models & Performance Evaluation",
            "subtitle": "Benchmarking Multiple Classification Algorithms",
            "bullets": [
                "Algorithms Tested: Baseline Logistic Regression, Decision Tree, Random Forest, Gradient Boosting.",
                "Baseline Performance (Test Set):",
                "  • Logistic Regression: Accuracy = 86.3%, F1-Score = 0.777, ROC-AUC = 0.926",
                "  • Decision Tree: Accuracy = 82.7%, F1-Score = 0.701, ROC-AUC = 0.789",
                "  • Random Forest: Accuracy = 87.7%, F1-Score = 0.791, ROC-AUC = 0.942",
                "Baseline Winner: Baseline Random Forest demonstrated strongest initial performance."
            ]
        },
        {
            "title": "6. Hyperparameter Optimization & Sensitivity Analysis",
            "subtitle": "5-Fold GridSearchCV Tuning & Parameter Sensitivity",
            "bullets": [
                "Optimization Method: 5-Fold GridSearchCV targeting F1-Score optimization.",
                "Random Forest Best Parameters: n_estimators = 300, max_depth = None, min_samples_split = 2, min_samples_leaf = 1, class_weight = 'balanced'.",
                "Decision Tree Best Parameters: max_depth = 5, criterion = 'entropy', class_weight = 'balanced'.",
                "Sensitivity Finding: Class weighting ('balanced') significantly boosted recall without sacrificing precision."
            ]
        },
        {
            "title": "7. Final Model Evaluation & Baseline Comparison",
            "subtitle": "Optimized Random Forest vs Baseline Models",
            "bullets": [
                "Final Evaluation Metrics (Test Set):",
                "  • Accuracy: 88.67% (+1.00% over baseline RF)",
                "  • Precision: 82.35% (+1.40% over baseline RF)",
                "  • Recall: 79.55% (+2.28% over baseline RF)",
                "  • F1-Score: 0.8092 (+1.85% over baseline RF)",
                "  • ROC-AUC: 0.9493 (+0.74% over baseline RF)",
                "Conclusion: Hyperparameter optimization yielded superior balance between Precision and Recall."
            ]
        },
        {
            "title": "8. Feature Importance Analysis",
            "subtitle": "Top Influential Variables Driving Purchase Predictions",
            "bullets": [
                "Top 5 Influential Features (Gini Importance):",
                "  1. CartItems (34.09%) - Strongest intent signal; direct proxy for carting decision.",
                "  2. TimeOnSite (16.34%) - Duration of browsing active engagement.",
                "  3. PreviousPurchases (10.62%) - Customer loyalty and past brand familiarity.",
                "  4. PagesViewed (9.12%) - Depth of product discovery.",
                "  5. DaysSinceLastVisit (6.18%) - Recency & churn indicator."
            ]
        },
        {
            "title": "9. Actionable Business Recommendations",
            "subtitle": "Translating Machine Learning Insights into Growth Strategies",
            "bullets": [
                "1. Automated Cart Recovery: Trigger exit popups and dynamic email drip for cart items >= 1.",
                "2. Dynamic Customer Segmentation: Allocate ad budget to High/Medium risk tiers; suppress Low tier.",
                "3. Customer Retention Drip: Launch automated 14-day & 30-day re-engagement offers.",
                "4. High-Browsing Assistance: Prompt live-chat help for users spending >4 min without carting.",
                "5. Dynamic Threshold Strategy: Lower threshold to 0.40 during peak holidays to maximize buyer capture."
            ]
        },
        {
            "title": "10. Project Limitations & Next Steps",
            "subtitle": "Future Enhancements & Production Roadmap",
            "bullets": [
                "Limitations:",
                "  • Feature Importance reflects correlation, not direct causality.",
                "  • Static offline dataset requires real-time streaming feature integration.",
                "Next Steps:",
                "  • Conduct randomized A/B testing on cart recovery triggers.",
                "  • Deploy model as REST API microservice using FastAPI & Docker.",
                "  • Implement SHAP value explanations for real-time model interpretability."
            ]
        }
    ]

    for slide_data in slides_content:
        slide = prs.slides.add_slide(blank_layout)
        
        # Add Header Background banner
        banner = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.333), Inches(1.3) # MSO_SHAPE.RECTANGLE
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = PRIMARY_COLOR
        banner.line.color.rgb = PRIMARY_COLOR

        # Title text
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.6))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data["title"]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Subtitle text
        tx_box_sub = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.4))
        tf_sub = tx_box_sub.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = slide_data["subtitle"]
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = ACCENT_COLOR

        # Content text box
        tx_box_body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.3))
        tf_body = tx_box_body.text_frame
        tf_body.word_wrap = True

        for i, bullet in enumerate(slide_data["bullets"]):
            p_b = tf_body.add_paragraph() if i > 0 else tf_body.paragraphs[0]
            p_b.text = bullet
            p_b.font.size = Pt(16)
            p_b.font.color.rgb = TEXT_DARK
            p_b.space_after = Pt(12)
            if bullet.startswith("  "):
                p_b.level = 1
                p_b.font.size = Pt(14)

    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    prs.save(output_pptx_path)
    print(f"[Presentation Generator] Successfully created presentation at '{output_pptx_path}'.")

if __name__ == "__main__":
    base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "presentation"))
    path1 = os.path.join(base_dir, "mini_project_5_presentation.pptx")
    create_presentation(path1)
