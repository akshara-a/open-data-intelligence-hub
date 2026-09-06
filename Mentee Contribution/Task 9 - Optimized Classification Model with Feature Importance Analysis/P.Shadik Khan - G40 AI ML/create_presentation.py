from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    slide.shapes.title.text = title

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for i, item in enumerate(content):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = item
        p.font.size = Pt(20)

# Slide 1
add_slide(
    "Purchase Prediction Model",
    [
        "Task 9 - Optimized Classification Model with Feature Importance Analysis",
        "Student: P.Shadik Khan - G40 AI ML",
        "Objective: Predict customers likely to make a purchase using machine learning."
    ]
)

# Slide 2
add_slide(
    "Business Problem",
    [
        "E-commerce companies need to identify customers who are likely to purchase.",
        "The prediction model helps prioritize customers for targeted marketing campaigns.",
        "The business objective is to improve customer targeting and conversion opportunities.",
        "This is a binary classification problem: Purchase vs No Purchase."
    ]
)

# Slide 3
add_slide(
    "Dataset and Exploratory Analysis",
    [
        "Dataset contains customer demographic, browsing and engagement information.",
        "Important variables include TimeOnSite, PagesViewed, ProductsViewed and CartItems.",
        "Exploratory analysis was performed using purchase distribution and feature relationships.",
        "The target distribution showed class imbalance, making accuracy alone insufficient."
    ]
)

# Slide 4
add_slide(
    "Preprocessing and Leakage Prevention",
    [
        "Data was split into training and unseen test data before model optimization.",
        "Numerical and categorical features were processed using a machine learning pipeline.",
        "Appropriate encoding and scaling were applied.",
        "The test dataset was not used during hyperparameter selection.",
        "A fixed random state was used for reproducibility."
    ]
)

# Slide 5
add_slide(
    "Models Compared",
    [
        "Logistic Regression - Baseline",
        "Decision Tree - Baseline",
        "Random Forest - Baseline",
        "Optimized Random Forest",
        "Models were evaluated using Accuracy, Precision, Recall, F1 Score and ROC-AUC."
    ]
)

# Slide 6
add_slide(
    "Hyperparameter Optimization",
    [
        "Model optimized: Random Forest",
        "Search method: GridSearchCV",
        "Cross-validation: 5-fold",
        "Optimization metric: F1 Score",
        "Best CV F1 Score: 0.0755",
        "Best model parameters included max_depth=8 and n_estimators=100."
    ]
)

# Slide 7
add_slide(
    "Model Performance and Threshold Analysis",
    [
        "Optimized Random Forest Accuracy: 0.8575",
        "Default threshold 0.50 produced no positive purchase predictions.",
        "Threshold 0.30 achieved Recall: 0.5319",
        "Threshold 0.30 achieved F1 Score: 0.2183",
        "This demonstrates why threshold tuning is important for imbalanced classification problems."
    ]
)

# Slide 8
add_slide(
    "Feature Importance",
    [
        "Top predictive features:",
        "1. TimeOnSite - 0.1323",
        "2. DaysSinceLastVisit - 0.1287",
        "3. ReviewScoreViewed - 0.0889",
        "4. AverageOrderValue - 0.0883",
        "5. Age - 0.0790",
        "Customer engagement and recency were important predictive signals.",
        "Feature importance indicates association and does not prove causation."
    ]
)

# Slide 9
add_slide(
    "Business Recommendations and Conclusion",
    [
        "Prioritize highly engaged customers for personalized campaigns.",
        "Use re-engagement strategies for customers who have not visited recently.",
        "Improve product reviews and product information visibility.",
        "Target customers who add products to their cart but do not purchase.",
        "Use threshold tuning based on campaign costs and business capacity.",
        "Future work: improve positive-class detection using class weighting or resampling."
    ]
)

output_file = r"Mentee Contribution\Task 9 - Optimized Classification Model with Feature Importance Analysis\P.Shadik Khan - G40 AI ML\presentation\mini_project_5_presentation.pptx"

prs.save(output_file)

print("Presentation created successfully!")
print(output_file)